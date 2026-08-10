# -*- coding: utf-8 -*-
"""Rendu de chaque type de slide, en s'appuyant sur design.py (moteur generique)."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import design as d

CONTENT_TOP = Inches(1.55)
CONTENT_LEFT = d.MARGIN_X
CONTENT_W = d.SLIDE_W - 2 * d.MARGIN_X


def _bullets_block(slide, x, y, w, h, items, size=17, color=d.TEXT, gap=8):
    _, tf = d.textbox(slide, x, y, w, h)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        d._no_bullet(p)
        r1 = p.add_run()
        d.set_run(r1, "›  ", size=size, color=d.ACCENT, bold=True)
        r2 = p.add_run()
        d.set_run(r2, it, size=size, color=color)


def render_cover(prs, e):
    slide = d.add_slide(prs)
    d.rect(slide, 0, 0, d.SLIDE_W, Inches(0.14), color=d.ACCENT)
    d.rect(slide, 0, d.SLIDE_H - Inches(0.14), d.SLIDE_W, Inches(0.14), color=d.ACCENT_2)
    _, tfk = d.textbox(slide, d.MARGIN_X, Inches(1.3), Inches(9), Inches(0.5))
    d.add_para(tfk, e["kicker"], size=16, color=d.ACCENT_2, bold=True, space_after=0)
    _, tft = d.textbox(slide, d.MARGIN_X, Inches(1.9), Inches(11.5), Inches(2.0))
    d.add_para(tft, e["title"], size=76, color=d.TEXT, bold=True, font=d.FONT_TITLE, space_after=0, line_spacing=1.0)
    _, tfs = d.textbox(slide, d.MARGIN_X, Inches(3.55), Inches(11), Inches(0.8))
    d.add_para(tfs, e["subtitle"], size=30, color=d.ACCENT, bold=True, font=d.FONT_TITLE, space_after=0)
    d.rect(slide, d.MARGIN_X, Inches(4.45), Inches(3.2), Pt(2), color=d.LINE_SOFT)
    _, tftag = d.textbox(slide, d.MARGIN_X, Inches(4.65), Inches(9.5), Inches(0.9))
    d.add_para(tftag, e["tagline"], size=17, color=d.TEXT_DIM, italic=True, space_after=0)
    _, tfm = d.textbox(slide, d.MARGIN_X, d.SLIDE_H - Inches(1.55), Inches(9), Inches(1.2))
    for m in e["meta"]:
        d.add_para(tfm, m, size=13, color=d.TEXT_FAINT, space_after=4)
    return slide


def render_section(prs, e):
    slide = d.add_slide(prs)
    d.rect(slide, 0, 0, Inches(0.22), d.SLIDE_H, color=d.ACCENT)
    _, tfk = d.textbox(slide, Inches(1.0), Inches(2.7), Inches(11), Inches(0.6))
    d.add_para(tfk, e["module_no"], size=20, color=d.ACCENT, bold=True, space_after=0)
    _, tft = d.textbox(slide, Inches(1.0), Inches(3.25), Inches(11.3), Inches(1.6))
    d.add_para(tft, e["title"], size=44, color=d.TEXT, bold=True, font=d.FONT_TITLE, space_after=0, line_spacing=1.02)
    if e.get("subtitle"):
        _, tfs = d.textbox(slide, Inches(1.0), Inches(4.35), Inches(10.8), Inches(1.0))
        d.add_para(tfs, e["subtitle"], size=18, color=d.TEXT_DIM, space_after=0)
    return slide


def _header(slide, e):
    d.kicker_title(slide, e.get("kicker", ""), e["title"])


def render_bullets(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP
    if e.get("definition"):
        box = d.rect(slide, CONTENT_LEFT, y, CONTENT_W, Inches(0.85), color=d.BG_PANEL, line_color=d.ACCENT, line_w=Pt(1))
        stripe = d.rect(slide, CONTENT_LEFT, y, Inches(0.08), Inches(0.85), color=d.ACCENT)
        _, tfd = d.textbox(slide, CONTENT_LEFT + Inches(0.3), y, CONTENT_W - Inches(0.5), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
        d.add_para(tfd, e["definition"], size=16, color=d.TEXT, italic=True, space_after=0)
        y += Inches(1.1)
    _bullets_block(slide, CONTENT_LEFT, y, CONTENT_W, d.SLIDE_H - y - Inches(0.7), e["bullets"])
    return slide


def render_code(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP
    h = d.SLIDE_H - y - Inches(0.75)
    if e.get("caption"):
        h -= Inches(0.55)
    d.code_box(slide, CONTENT_LEFT, y, CONTENT_W, h, e["code"], lang=e.get("lang", "java"), size=e.get("size", 17))
    if e.get("caption"):
        _, tfc = d.textbox(slide, CONTENT_LEFT, y + h + Inches(0.12), CONTENT_W, Inches(0.45))
        d.add_para(tfc, e["caption"], size=14, color=d.ACCENT_2, italic=True, space_after=0)
    return slide


def render_code_bullets(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP
    h = d.SLIDE_H - y - Inches(0.7)
    half = (CONTENT_W - Inches(0.4)) / 2
    d.code_box(slide, CONTENT_LEFT, y, half, h, e["code"], lang=e.get("lang", "java"), size=15)
    _bullets_block(slide, CONTENT_LEFT + half + Inches(0.4), y + Inches(0.1), half, h, e["bullets"], size=15)
    return slide


def render_diagram_flow(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    boxes = e["boxes"]
    n = len(boxes)
    y = CONTENT_TOP + Inches(0.15)
    avail_h = d.SLIDE_H - y - Inches(0.7)
    gap = Inches(0.3)
    box_h = min(Inches(0.85), (avail_h - gap * (n - 1)) / n)
    w = Inches(8.6)
    x = d.SLIDE_W / 2 - w / 2
    d.flow_diagram(slide, x, y, w, box_h, boxes, gap=gap)
    return slide


def render_diagram_tree(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    d.tree_diagram(slide, e["root"], e["children"], CONTENT_LEFT, CONTENT_TOP + Inches(0.3), CONTENT_W)
    return slide


def render_definitions_trio(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    items = e["items"]
    n = len(items)
    gap = Inches(0.35)
    w = (CONTENT_W - gap * (n - 1)) / n
    y = CONTENT_TOP + Inches(0.2)
    h = Inches(3.6)
    for i, (term, desc) in enumerate(items):
        x = CONTENT_LEFT + i * (w + gap)
        d.rect(slide, x, y, w, h, color=d.BG_PANEL, line_color=d.LINE_SOFT, line_w=Pt(1))
        d.rect(slide, x, y, w, Inches(0.08), color=d.ACCENT if i % 2 == 0 else d.ACCENT_2)
        _, tft = d.textbox(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.9))
        d.add_para(tft, term, size=19, color=d.TEXT, bold=True, font=d.FONT_TITLE, space_after=0, line_spacing=1.0)
        _, tfd = d.textbox(slide, x + Inches(0.25), y + Inches(1.15), w - Inches(0.5), h - Inches(1.4))
        d.add_para(tfd, desc, size=14, color=d.TEXT_DIM, space_after=0, line_spacing=1.2)
    return slide


def render_table(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    headers = e["headers"]
    rows = e["rows"]
    ncols = len(headers)
    nrows = len(rows) + 1
    y = CONTENT_TOP + Inches(0.1)
    h = min(Inches(0.62) * nrows, d.SLIDE_H - y - Inches(0.7))
    w = CONTENT_W
    gshape = slide.shapes.add_table(nrows, ncols, CONTENT_LEFT, y, w, h)
    table = gshape.table
    for c in range(ncols):
        table.columns[c].width = int(w / ncols)
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = d.ACCENT
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4); cell.margin_left = Pt(10); cell.margin_right = Pt(10)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        d.set_run(r, htext, size=14, color=RGBColor(0x0B, 0x12, 0x20), bold=True)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = d.BG_PANEL if ri % 2 == 1 else d.BG
            cell.margin_top = Pt(4); cell.margin_bottom = Pt(4); cell.margin_left = Pt(10); cell.margin_right = Pt(10)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            weight = (c == 0)
            d.set_run(r, str(val), size=13.5, color=d.TEXT, bold=weight)
    return slide


def render_errors(prs, e):
    slide = d.add_slide(prs)
    d.kicker_title(slide, e.get("kicker", ""), e["title"], kicker_color=d.BAD)
    y = CONTENT_TOP
    _, tf = d.textbox(slide, CONTENT_LEFT, y, CONTENT_W, d.SLIDE_H - y - Inches(0.7))
    first = True
    for it in e["items"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        d._no_bullet(p)
        r1 = p.add_run()
        d.set_run(r1, "✖  ", size=17, color=d.BAD, bold=True)
        r2 = p.add_run()
        d.set_run(r2, it, size=17, color=d.TEXT)
    return slide


def render_exercise(prs, e):
    slide = d.add_slide(prs)
    d.kicker_title(slide, e.get("kicker", "À VOUS DE JOUER"), e["title"], kicker_color=d.GOOD)
    y = CONTENT_TOP + Inches(0.15)
    box = d.rect(slide, CONTENT_LEFT, y, CONTENT_W, Inches(2.2), color=d.BG_PANEL, line_color=d.GOOD, line_w=Pt(1.25))
    stripe = d.rect(slide, CONTENT_LEFT, y, Inches(0.1), Inches(2.2), color=d.GOOD)
    _, tf = d.textbox(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.3), CONTENT_W - Inches(0.8), Inches(1.7))
    d.add_para(tf, e["prompt"], size=19, color=d.TEXT, space_after=0, line_spacing=1.3)
    if e.get("requirements"):
        y2 = y + Inches(2.45)
        _bullets_block(slide, CONTENT_LEFT, y2, CONTENT_W, Inches(1.8), e["requirements"], size=15)
    return slide


def render_summary(prs, e):
    slide = d.add_slide(prs)
    d.kicker_title(slide, e.get("kicker", ""), e.get("title", "À RETENIR"), kicker_color=d.GOLD)
    y = CONTENT_TOP + Inches(0.1)
    box = d.rect(slide, CONTENT_LEFT, y, CONTENT_W, d.SLIDE_H - y - Inches(0.6), color=d.BG_PANEL, line_color=d.GOLD, line_w=Pt(1))
    _bullets_block(slide, CONTENT_LEFT + Inches(0.4), y + Inches(0.35), CONTENT_W - Inches(0.8), d.SLIDE_H - y - Inches(1.1), e["bullets"], size=18, gap=14)
    return slide


def render_quiz(prs, e):
    slide = d.add_slide(prs)
    d.kicker_title(slide, e.get("kicker", ""), e.get("title", "QUIZ"), kicker_color=d.GOLD)
    y = CONTENT_TOP
    _, tf = d.textbox(slide, CONTENT_LEFT, y, CONTENT_W, d.SLIDE_H - y - Inches(0.7))
    first = True
    for i, q in enumerate(e["questions"], start=1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        p.line_spacing = 1.2
        d._no_bullet(p)
        r1 = p.add_run()
        d.set_run(r1, f"{i}.  ", size=16.5, color=d.GOLD, bold=True)
        r2 = p.add_run()
        d.set_run(r2, q, size=16.5, color=d.TEXT)
    return slide


def render_compare(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP + Inches(0.1)
    h = d.SLIDE_H - y - Inches(0.7)
    half = (CONTENT_W - Inches(0.4)) / 2
    for i, (title, items, accent) in enumerate([(e["left_title"], e["left_items"], d.ACCENT), (e["right_title"], e["right_items"], d.ACCENT_2)]):
        x = CONTENT_LEFT + i * (half + Inches(0.4))
        d.rect(slide, x, y, half, h, color=d.BG_PANEL, line_color=d.LINE_SOFT, line_w=Pt(1))
        d.rect(slide, x, y, half, Inches(0.55), color=accent)
        _, tft = d.textbox(slide, x, y, half, Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
        d.add_para(tft, title, size=17, color=RGBColor(0x0B, 0x12, 0x20), bold=True, align=PP_ALIGN.CENTER, space_after=0)
        _bullets_block(slide, x + Inches(0.3), y + Inches(0.85), half - Inches(0.6), h - Inches(1.1), items, size=14.5)
    return slide


def render_diagram_class_uml(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP + Inches(0.4)
    w = Inches(4.6)
    x1 = CONTENT_LEFT + Inches(0.3)
    x2 = d.SLIDE_W - d.MARGIN_X - Inches(0.3) - w
    h1 = d.uml_class_box(slide, x1, y, w, e["class_name"], attrs=e.get("attrs"), methods=e.get("methods"))
    h2 = d.uml_class_box(slide, x2, y, w, e["class2_name"], attrs=e.get("attrs2"), methods=e.get("methods2"))
    mid_y = y + min(h1, h2) / 2
    d.uml_link(slide, x1 + w, mid_y, x2, mid_y)
    if e.get("multiplicity"):
        _, tfm = d.textbox(slide, x1 + w + Inches(0.15), mid_y - Inches(0.55), x2 - x1 - w - Inches(0.3), Inches(0.4))
        d.add_para(tfm, e["multiplicity"], size=13, color=d.ACCENT_2, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    return slide


def render_question(prs, e):
    slide = d.add_slide(prs)
    panel = d.rect(slide, 0, 0, d.SLIDE_W, d.SLIDE_H, color=d.BG_PANEL)
    d._send_to_back(slide, panel)
    _, tfk = d.textbox(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(0.6))
    d.add_para(tfk, e.get("title", "QUESTION"), size=22, color=d.GOLD, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tft = d.textbox(slide, Inches(1.2), Inches(3.3), Inches(10.9), Inches(2.0))
    d.add_para(tft, e["text"], size=32, color=d.TEXT, bold=True, font=d.FONT_TITLE, align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.15)
    return slide


def render_answer(prs, e):
    slide = d.add_slide(prs)
    d.kicker_title(slide, "", e.get("title", "RÉPONSE"), kicker_color=d.GOOD, title_size=30)
    _bullets_block(slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.3), CONTENT_W, Inches(3.5), e["bullets"], size=19, gap=16)
    return slide


def render_tp(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    y = CONTENT_TOP
    _, tfo = d.textbox(slide, CONTENT_LEFT, y, CONTENT_W, Inches(0.7))
    d.add_para(tfo, e["objective"], size=17, color=d.ACCENT_2, italic=True, space_after=0, line_spacing=1.2)
    y += Inches(0.85)
    _bullets_block(slide, CONTENT_LEFT, y, CONTENT_W, Inches(2.6), e["requirements"], size=16.5)
    _, tfs = d.textbox(slide, CONTENT_LEFT, d.SLIDE_H - Inches(0.95), CONTENT_W, Inches(0.35))
    d.add_para(tfs, "Notions mobilisées : " + e.get("skills", ""), size=12.5, color=d.TEXT_FAINT, italic=True, space_after=0)
    return slide


def render_glossary(prs, e):
    slide = d.add_slide(prs)
    _header(slide, e)
    items = e["items"]
    mid = (len(items) + 1) // 2
    col_w = (CONTENT_W - Inches(0.5)) / 2
    for ci, chunk in enumerate([items[:mid], items[mid:]]):
        x = CONTENT_LEFT + ci * (col_w + Inches(0.5))
        _, tf = d.textbox(slide, x, CONTENT_TOP, col_w, d.SLIDE_H - CONTENT_TOP - Inches(0.7))
        first = True
        for term, desc in chunk:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10)
            p.line_spacing = 1.1
            d._no_bullet(p)
            r1 = p.add_run()
            d.set_run(r1, term + " — ", size=13.5, color=d.ACCENT_2, bold=True)
            r2 = p.add_run()
            d.set_run(r2, desc, size=13, color=d.TEXT_DIM)
    return slide


def render_closing(prs, e):
    slide = d.add_slide(prs)
    d.rect(slide, 0, Inches(3.2), d.SLIDE_W, Inches(0.06), color=d.ACCENT)
    _, tfk = d.textbox(slide, Inches(1.0), Inches(2.2), Inches(11), Inches(0.5))
    d.add_para(tfk, e.get("kicker", ""), size=15, color=d.ACCENT_2, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tft = d.textbox(slide, Inches(1.2), Inches(3.5), Inches(10.9), Inches(2.2))
    d.add_para(tft, e["quote"], size=30, color=d.TEXT, bold=True, italic=True, font=d.FONT_TITLE, align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.25)
    return slide


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "bullets": render_bullets,
    "code": render_code,
    "code_bullets": render_code_bullets,
    "diagram_flow": render_diagram_flow,
    "diagram_tree": render_diagram_tree,
    "diagram_class_uml": render_diagram_class_uml,
    "definitions_trio": render_definitions_trio,
    "table": render_table,
    "errors": render_errors,
    "exercise": render_exercise,
    "summary": render_summary,
    "quiz": render_quiz,
    "compare": render_compare,
    "question": render_question,
    "answer": render_answer,
    "tp": render_tp,
    "glossary": render_glossary,
    "closing": render_closing,
}
