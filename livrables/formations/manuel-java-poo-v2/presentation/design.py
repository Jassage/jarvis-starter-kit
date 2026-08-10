# -*- coding: utf-8 -*-
"""
Moteur de rendu (design system) pour la presentation
"JAVA - Programmation Orientee Objet".

Ce module ne contient AUCUN contenu pedagogique : uniquement les fonctions
de rendu reutilisables (cover, section, bullets, code, diagrammes, tableaux,
exercices, quiz, resume, erreurs frequentes, TP...).

Le contenu reel est dans les fichiers content_*.py (une liste SLIDES par
partie du cours), assembles par build.py.
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette --
BG        = RGBColor(0x0B, 0x12, 0x20)   # fond principal (navy tres sombre)
BG_PANEL  = RGBColor(0x12, 0x1C, 0x32)   # panneaux / cartes
BG_CODE   = RGBColor(0x0A, 0x10, 0x1D)   # fond des blocs de code
LINE_SOFT = RGBColor(0x24, 0x31, 0x4D)   # liserés / séparateurs

ACCENT    = RGBColor(0xF2, 0x8C, 0x28)   # orange Java (accent principal)
ACCENT_2  = RGBColor(0x35, 0xC7, 0xE0)   # cyan (accent secondaire)
GOOD      = RGBColor(0x33, 0xC9, 0x84)   # vert (exercices / corrections)
BAD       = RGBColor(0xE8, 0x5C, 0x5C)   # rouge (erreurs frequentes)
GOLD      = RGBColor(0xE8, 0xC5, 0x4A)   # or (quiz / mise en avant)

TEXT      = RGBColor(0xEC, 0xF0, 0xF7)   # texte principal
TEXT_DIM  = RGBColor(0x93, 0xA1, 0xBE)   # texte secondaire / muted
TEXT_FAINT= RGBColor(0x5C, 0x6A, 0x87)   # footer / mentions

CODE_TEXT    = RGBColor(0xE3, 0xE9, 0xF4)
CODE_KEYWORD = RGBColor(0xF2, 0x8C, 0x28)
CODE_TYPE    = RGBColor(0x35, 0xC7, 0xE0)
CODE_STRING  = RGBColor(0x33, 0xC9, 0x84)
CODE_COMMENT = RGBColor(0x5C, 0x6A, 0x87)
CODE_NUMBER  = RGBColor(0xE8, 0xC5, 0x4A)

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY  = "Segoe UI"
FONT_CODE  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.7)

# --------------------------------------------------------------- helpers --

def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # layout vierge


def add_slide(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    _bg(slide, BG)
    return slide


def _bg(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    _send_to_back(slide, rect)
    return rect


def _send_to_back(slide, shape):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def no_line(shape):
    shape.line.fill.background()


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.shadow.inherit = False


def _i(v):
    """Force un entier EMU. python-pptx/PowerPoint corrompt silencieusement le fichier
    (refus d'ouverture, sans message de reparation) si une coordonnee flottante
    (issue par ex. d'une division '/') est ecrite dans le XML."""
    return int(round(v))


def rect(slide, x, y, w, h, color=None, line_color=None, line_w=Pt(1), shape=MSO_SHAPE.RECTANGLE):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if color is None:
        s.fill.background()
    else:
        solid(s, color)
    if line_color is None:
        no_line(s)
    else:
        s.line.color.rgb = line_color
        s.line.width = line_w
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size=18, color=TEXT, bold=False, italic=False, font=FONT_BODY):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return run


def add_para(tf, text, size=18, color=TEXT, bold=False, italic=False, font=FONT_BODY,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0, bullet=False,
             line_spacing=None):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    if bullet:
        _set_bullet(p, color)
    else:
        _no_bullet(p)
    return p


def _no_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    buNone = pPr.makeelement(qn('a:buNone'), {})
    pPr.append(buNone)


def _set_bullet(p, color):
    pPr = p._p.get_or_add_pPr()
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '‣'})
    pPr.append(buFont)
    pPr.append(buChar)


def kicker_title(slide, kicker, title, y=Inches(0.55), title_size=32, kicker_color=ACCENT,
                  title_color=TEXT, x=MARGIN_X, w=None):
    if w is None:
        w = SLIDE_W - 2 * MARGIN_X
    if kicker:
        _, tf = textbox(slide, x, y, w, Inches(0.4))
        add_para(tf, kicker.upper(), size=13, color=kicker_color, bold=True, font=FONT_BODY, space_after=0)
        title_y = y + Inches(0.36)
    else:
        title_y = y
    _, tf2 = textbox(slide, x, title_y, w, Inches(1.0))
    add_para(tf2, title, size=title_size, color=title_color, bold=True, font=FONT_TITLE, space_after=0,
              line_spacing=1.0)
    bar = rect(slide, x, title_y + Inches(0.02), Inches(0.06), Inches(0.5), color=ACCENT)
    return title_y


def footer(slide, module_label, page_no, total=None):
    _, tf = textbox(slide, MARGIN_X, SLIDE_H - Inches(0.42), Inches(7), Inches(0.3))
    add_para(tf, "JAVA — PROGRAMMATION ORIENTÉE OBJET", size=9, color=TEXT_FAINT, font=FONT_BODY, space_after=0)
    _, tf2 = textbox(slide, SLIDE_W - Inches(4.2), SLIDE_H - Inches(0.42), Inches(3.5), Inches(0.3))
    label = module_label if module_label else ""
    txt = f"{label}"
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, txt, size=9, color=TEXT_FAINT)
    if page_no is not None:
        _, tf3 = textbox(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.42), Inches(0.5), Inches(0.3))
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        set_run(r3, str(page_no), size=9, color=TEXT_FAINT)
    rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Pt(1), color=LINE_SOFT)


def set_notes(slide, text):
    if not text:
        return
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ------------------------------------------------------------ code coloring --

JAVA_KEYWORDS = set("""public private protected static final void class interface extends implements
new return if else for while do switch case break continue try catch finally throw throws
this super import package abstract enum default synchronized volatile transient instanceof
true false null int double float long short byte char boolean String var record sealed permits
""".split())

JAVA_TYPES = set("""String Integer Double Boolean List ArrayList LinkedList Map HashMap HashSet Set
Optional Stream Scanner Connection Statement PreparedStatement ResultSet DriverManager
SQLException Exception RuntimeException IllegalArgumentException ArithmeticException Object
System Math Files Path Paths File BufferedReader FileReader FileWriter Comparator Function
Predicate Consumer Supplier Collectors JOptionPane
""".split())

TOKEN_RE = re.compile(
    r'(?P<comment>//[^\n]*)'
    r'|(?P<string>"(?:[^"\\]|\\.)*")'
    r'|(?P<char>\'(?:[^\'\\]|\\.)*\')'
    r'|(?P<number>\b\d+\.?\d*\b)'
    r'|(?P<word>[A-Za-z_][A-Za-z0-9_]*)'
    r'|(?P<other>.)'
    , re.DOTALL)


def tokenize_java_line(line):
    """Retourne une liste de (texte, couleur) pour une ligne de code Java."""
    out = []
    for m in TOKEN_RE.finditer(line):
        kind = m.lastgroup
        text = m.group()
        if kind == 'comment':
            out.append((text, CODE_COMMENT))
        elif kind == 'string':
            out.append((text, CODE_STRING))
        elif kind == 'char':
            out.append((text, CODE_STRING))
        elif kind == 'number':
            out.append((text, CODE_NUMBER))
        elif kind == 'word':
            if text in JAVA_KEYWORDS:
                out.append((text, CODE_KEYWORD))
            elif text in JAVA_TYPES:
                out.append((text, CODE_TYPE))
            else:
                out.append((text, CODE_TEXT))
        else:
            out.append((text, CODE_TEXT))
    return out


SQL_KEYWORDS = set("""SELECT FROM WHERE INSERT INTO VALUES UPDATE SET DELETE CREATE TABLE
PRIMARY KEY FOREIGN REFERENCES NOT NULL DEFAULT AUTO_INCREMENT VARCHAR INT DOUBLE DECIMAL
DATE BOOLEAN JOIN INNER LEFT RIGHT ON GROUP BY ORDER HAVING AND OR IN LIKE AS ALTER DROP
UNIQUE CONSTRAINT ENGINE CHARSET LIMIT COUNT SUM AVG MAX MIN DISTINCT""".split())


def tokenize_sql_line(line):
    out = []
    for m in TOKEN_RE.finditer(line):
        kind = m.lastgroup
        text = m.group()
        if kind == 'comment':
            out.append((text, CODE_COMMENT))
        elif kind == 'string':
            out.append((text, CODE_STRING))
        elif kind == 'number':
            out.append((text, CODE_NUMBER))
        elif kind == 'word':
            if text.upper() in SQL_KEYWORDS:
                out.append((text, CODE_KEYWORD))
            else:
                out.append((text, CODE_TEXT))
        else:
            out.append((text, CODE_TEXT))
    return out


def code_box(slide, x, y, w, h, code, lang='java', size=15, label=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    panel = rect(slide, x, y, w, h, color=BG_CODE, line_color=LINE_SOFT, line_w=Pt(1))
    if label:
        tag = rect(slide, x, y - Inches(0.02), Inches(1.7), Inches(0.32), color=BG_PANEL)
        tag.top = y - Inches(0.16)
        tag.left = x + Inches(0.15)
        _, tftag = textbox(slide, tag.left + Inches(0.1), tag.top + Inches(0.02), Inches(1.5), Inches(0.28))
        add_para(tftag, label, size=10, color=ACCENT_2, bold=True, font=FONT_CODE, space_after=0)
    _, tf = textbox(slide, x + Inches(0.28), y + Inches(0.2), w - Inches(0.56), h - Inches(0.4))
    lines = code.strip("\n").split("\n")
    tokenize = tokenize_java_line if lang == 'java' else tokenize_sql_line
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        p.line_spacing = 1.15
        _no_bullet(p)
        if line.strip() == "":
            r = p.add_run()
            set_run(r, " ", size=size, color=CODE_TEXT, font=FONT_CODE)
            continue
        for text, color in tokenize(line):
            r = p.add_run()
            set_run(r, text, size=size, color=color, font=FONT_CODE)
    return panel


# ------------------------------------------------------------ diagrams --

def flow_diagram(slide, x, y, w, box_h, labels, sub=None, vertical=True, gap=Inches(0.35),
                  box_color=BG_PANEL, text_color=TEXT, accent=ACCENT, arrow_color=None):
    """Empile des boites reliees par des fleches (verticales par defaut)."""
    if arrow_color is None:
        arrow_color = accent
    n = len(labels)
    cy = y
    for i, label in enumerate(labels):
        b = rect(slide, x, cy, w, box_h, color=box_color, line_color=LINE_SOFT, line_w=Pt(1))
        stripe = rect(slide, x, cy, Inches(0.09), box_h, color=accent)
        _, tf = textbox(slide, x + Inches(0.3), cy, w - Inches(0.5), box_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, label, size=17, color=text_color, bold=True, font=FONT_BODY, align=PP_ALIGN.LEFT, space_after=0)
        if sub and sub[i]:
            _, tf2 = textbox(slide, x + Inches(0.3), cy + box_h - Inches(0.02), w - Inches(0.5), Inches(0.3),
                              anchor=MSO_ANCHOR.TOP)
        cy += box_h
        if i < n - 1:
            ah = gap
            arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, _i(x + w / 2 - Inches(0.14)), _i(cy + Inches(0.04)),
                                          Inches(0.28), _i(ah - Inches(0.08)))
            solid(arr, arrow_color)
            no_line(arr)
            cy += ah
    return cy


def hflow_diagram(slide, x, y, box_w, h, labels, gap=Inches(0.3), box_color=BG_PANEL, text_color=TEXT,
                   accent=ACCENT):
    n = len(labels)
    cx = x
    for i, label in enumerate(labels):
        b = rect(slide, cx, y, box_w, h, color=box_color, line_color=LINE_SOFT, line_w=Pt(1))
        stripe = rect(slide, cx, y, box_w, Inches(0.08), color=accent)
        _, tf = textbox(slide, cx + Inches(0.15), y, box_w - Inches(0.3), h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, label, size=15, color=text_color, bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER, space_after=0)
        cx += box_w
        if i < n - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, _i(cx + Inches(0.03)), _i(y + h / 2 - Inches(0.12)),
                                          _i(gap - Inches(0.06)), Inches(0.24))
            solid(arr, accent)
            no_line(arr)
            cx += gap
    return cx


def tree_diagram(slide, root_label, children, x, y, w, root_color=ACCENT_2, child_color=BG_PANEL):
    """Boite racine en haut, N enfants en dessous relies par des traits."""
    root_w = Inches(2.6)
    root_h = Inches(0.6)
    root_x = _i(x + w / 2 - root_w / 2)
    rbox = rect(slide, root_x, y, root_w, root_h, color=BG_PANEL, line_color=root_color, line_w=Pt(1.5))
    _, tf = textbox(slide, root_x, y, root_w, root_h, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, root_label, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER, space_after=0)

    n = len(children)
    child_w = Inches(2.4)
    child_h = Inches(0.55)
    total_w = n * child_w + (n - 1) * Inches(0.5)
    start_x = _i(x + w / 2 - total_w / 2)
    mid_y = _i(y + root_h + Inches(0.35))
    trunk = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x + w / 2), y + root_h, _i(x + w / 2), mid_y)
    trunk.line.color.rgb = LINE_SOFT
    trunk.line.width = Pt(2)
    if n > 1:
        branch = rect(slide, _i(start_x + child_w / 2), mid_y, total_w - child_w, Pt(2), color=LINE_SOFT)
    cx = start_x
    for label in children:
        vconn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(cx + child_w / 2), mid_y,
                                            _i(cx + child_w / 2), _i(mid_y + Inches(0.25)))
        vconn.line.color.rgb = LINE_SOFT
        vconn.line.width = Pt(2)
        cbox = rect(slide, cx, mid_y + Inches(0.25), child_w, child_h, color=BG_PANEL,
                    line_color=LINE_SOFT, line_w=Pt(1))
        stripe = rect(slide, cx, mid_y + Inches(0.25), Inches(0.08), child_h, color=ACCENT)
        _, tf2 = textbox(slide, cx, mid_y + Inches(0.25), child_w, child_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf2, label, size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        cx += child_w + Inches(0.5)
    return mid_y + Inches(0.25) + child_h


def uml_class_box(slide, x, y, w, name, attrs=None, methods=None, abstract=False, stereotype=None,
                   name_color=None):
    attrs = attrs or []
    methods = methods or []
    line_h = Inches(0.28)
    header_h = Inches(0.55) if not stereotype else Inches(0.75)
    attrs_h = max(Inches(0.15), line_h * max(1, len(attrs)))
    methods_h = max(Inches(0.15), line_h * max(1, len(methods)))
    total_h = header_h + attrs_h + methods_h
    color = name_color or (ACCENT_2)
    outer = rect(slide, x, y, w, total_h, color=BG_PANEL, line_color=color, line_w=Pt(1.5))
    cy = y
    _, tfh = textbox(slide, x, cy, w, header_h, anchor=MSO_ANCHOR.MIDDLE)
    if stereotype:
        add_para(tfh, f"«{stereotype}»", size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER, space_after=0)
        p2 = tfh.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        style = "Segoe UI"
        set_run(r2, name, size=16, color=TEXT, bold=True, italic=abstract, font=FONT_TITLE)
        _no_bullet(p2)
    else:
        add_para(tfh, name, size=17, color=TEXT, bold=True, italic=abstract, align=PP_ALIGN.CENTER,
                  font=FONT_TITLE, space_after=0)
    cy += header_h
    rect(slide, x, cy, w, Pt(1.5), color=color)
    _, tfa = textbox(slide, x + Inches(0.15), cy + Inches(0.06), w - Inches(0.3), attrs_h)
    if attrs:
        for a in attrs:
            add_para(tfa, a, size=12.5, color=TEXT_DIM, font=FONT_CODE, space_after=2)
    else:
        add_para(tfa, " ", size=12, color=TEXT_DIM, space_after=0)
    cy += attrs_h
    rect(slide, x, cy, w, Pt(1.5), color=color)
    _, tfm = textbox(slide, x + Inches(0.15), cy + Inches(0.06), w - Inches(0.3), methods_h)
    if methods:
        for mtext in methods:
            add_para(tfm, mtext, size=12.5, color=TEXT, font=FONT_CODE, space_after=2)
    else:
        add_para(tfm, " ", size=12, color=TEXT, space_after=0)
    return total_h


def uml_link(slide, x1, y1, x2, y2, dashed=False, color=LINE_SOFT):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.75)
    if dashed:
        ln = conn.line._get_or_add_ln()
        dash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(dash)
    return conn
